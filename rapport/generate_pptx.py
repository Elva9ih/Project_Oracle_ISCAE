"""Generate PowerPoint presentation for Projet Oracle."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(__file__)
SCREENSHOT_DIR = os.path.join(SCRIPT_DIR, "screenshots")
LOGO_PATH = os.path.join(os.path.dirname(SCRIPT_DIR), "logo-iscae.png")
OUTPUT = os.path.join(os.path.dirname(SCRIPT_DIR), "Presentation_Projet_Oracle.pptx")

# Colors
GREEN_DARK = RGBColor(0x0D, 0x5C, 0x28)
GREEN_MAIN = RGBColor(0x1A, 0x7A, 0x3A)
GREEN_LIGHT = RGBColor(0x4C, 0xAF, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
GRAY_MED = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xE6, 0x51, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=GRAY_DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, items, left, top, width, height, font_size=16, color=GRAY_DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0

def add_footer(slide, number):
    # Green bottom bar
    add_shape(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), GREEN_MAIN)
    add_text_box(slide, Inches(0.5), Inches(7.12), Inches(6), Inches(0.35),
                 "Projet Oracle - ISCAE Nouakchott - 2025/2026", 10, WHITE)
    add_text_box(slide, Inches(11), Inches(7.12), Inches(2), Inches(0.35),
                 str(number), 10, WHITE, alignment=PP_ALIGN.RIGHT)

def add_screenshot(slide, filename, left, top, width, height=None):
    path = os.path.join(SCREENSHOT_DIR, filename)
    if os.path.exists(path):
        if height:
            slide.shapes.add_picture(path, left, top, width, height)
        else:
            slide.shapes.add_picture(path, left, top, width=width)

page = 0

# ============================================================
# SLIDE 1: Cover
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)

# Top green bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), GREEN_MAIN)

# Logo
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(5.6), Inches(0.5), Inches(2), Inches(2))

# Title
add_text_box(slide, Inches(1), Inches(2.8), Inches(11.333), Inches(1),
             "Oracle dans un Écosystème Applicatif", 36, GREEN_DARK, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1), Inches(3.8), Inches(11.333), Inches(0.6),
             "Gestion des Absences Étudiants — Django + Oracle 18c", 22, GRAY_MED, alignment=PP_ALIGN.CENTER)

# Green divider
add_shape(slide, Inches(5.5), Inches(4.5), Inches(2.333), Inches(0.04), GREEN_MAIN)

# Students
add_text_box(slide, Inches(1), Inches(4.9), Inches(11.333), Inches(0.5),
             "Réalisé par :", 14, GRAY_MED, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.3), Inches(11.333), Inches(0.5),
             "IP22240 — Cherifa Mohamed Yeslem Bellahi", 18, GRAY_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.75), Inches(11.333), Inches(0.5),
             "IP22237 — El Vagih Ahmed Zeine", 18, GRAY_DARK, bold=True, alignment=PP_ALIGN.CENTER)

# Institution
add_text_box(slide, Inches(1), Inches(6.4), Inches(11.333), Inches(0.5),
             "ISCAE Nouakchott — Master MPIAG — Année universitaire 2025-2026", 14, GRAY_MED, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), GREEN_MAIN)


# ============================================================
# SLIDE 2: Plan
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Plan de la Présentation", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

items = [
    "1.  Contexte et objectifs du projet",
    "2.  Architecture technique (Oracle + Django)",
    "3.  Modèle de données (schéma relationnel)",
    "4.  PL/SQL : Fonctions, Procédures, Packages",
    "5.  Triggers et Audit automatique",
    "6.  Vues et Vues Matérialisées",
    "7.  Optimisation et EXPLAIN PLAN",
    "8.  Application Django — Démonstration",
    "9.  Gestion des rôles (Admin, Enseignant, Étudiant)",
    "10. Export (PDF, Excel) et Sauvegarde (Data Pump)",
    "11. Conclusion et perspectives",
]
add_bullet_slide(slide, items, Inches(1.5), Inches(1.6), Inches(10), Inches(5.5), 20, GRAY_DARK)
add_footer(slide, page)


# ============================================================
# SLIDE 3: Contexte
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "1. Contexte et Objectifs", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

# Left column - Contexte
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
             "Contexte", 22, GREEN_MAIN, bold=True)
items = [
    "Projet du module « Oracle dans un Écosystème Applicatif »",
    "Sujet : Gestion des Absences Étudiants",
    "Base de données Oracle 18c XE",
    "Application web Django (Python 3.11)",
    "Répartition : 75% Oracle / 25% Django",
]
add_bullet_slide(slide, items, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3), 16, GRAY_DARK)

# Right column - Objectifs
add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
             "Objectifs", 22, GREEN_MAIN, bold=True)
items = [
    "Exploiter les fonctionnalités avancées d'Oracle",
    "PL/SQL : packages, fonctions, procédures",
    "Triggers d'audit et de cohérence",
    "Vues matérialisées pour la performance",
    "Intégration Django ↔ Oracle (cx_Oracle)",
    "Gestion des rôles et permissions",
]
add_bullet_slide(slide, items, Inches(7), Inches(2.1), Inches(5.5), Inches(3.5), 16, GRAY_DARK)

# Divider
add_shape(slide, Inches(6.5), Inches(1.5), Inches(0.03), Inches(4.5), GREEN_LIGHT)

add_footer(slide, page)


# ============================================================
# SLIDE 4: Architecture
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "2. Architecture Technique", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

# Architecture boxes
# Oracle box
box = add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), GRAY_LIGHT)
add_text_box(slide, Inches(1), Inches(1.7), Inches(5), Inches(0.5),
             "Oracle 18c XE (75%)", 22, GREEN_DARK, bold=True, alignment=PP_ALIGN.CENTER)
items = [
    "12 tables avec contraintes complètes",
    "23+ index (FK, fonctionnels)",
    "2 packages PL/SQL",
    "5 fonctions + 6 procédures",
    "8 triggers (audit + métier)",
    "5 vues + 4 vues matérialisées",
    "120 000+ lignes de données",
    "Transactions (COMMIT, ROLLBACK, SAVEPOINT)",
    "Export Data Pump",
]
add_bullet_slide(slide, items, Inches(1.2), Inches(2.3), Inches(5), Inches(4), 15, GRAY_DARK)

# Django box
box = add_shape(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5), GRAY_LIGHT)
add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.5),
             "Django 3.2 + Python 3.11 (25%)", 22, GREEN_DARK, bold=True, alignment=PP_ALIGN.CENTER)
items = [
    "CRUD complet (étudiants, enseignants, séances)",
    "Marquage d'absences (appel PL/SQL)",
    "Gestion des justificatifs",
    "Tableau de bord avec Chart.js",
    "3 rôles : Admin, Enseignant, Étudiant",
    "Export PDF (ReportLab) et Excel (openpyxl)",
    "Journal d'audit (depuis triggers Oracle)",
    "Interface Bootstrap 5 responsive",
    "Connexion via oracledb / cx_Oracle",
]
add_bullet_slide(slide, items, Inches(7.4), Inches(2.3), Inches(5), Inches(4), 15, GRAY_DARK)

# Arrow between
add_text_box(slide, Inches(6.1), Inches(3.5), Inches(1.2), Inches(1),
             "⟷", 40, GREEN_MAIN, bold=True, alignment=PP_ALIGN.CENTER)

add_footer(slide, page)


# ============================================================
# SLIDE 5: Modèle de données
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "3. Modèle de Données", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

# Tables list in 3 columns
tables_col1 = [
    "FILIERES — Filières d'études",
    "GROUPES — Groupes par filière",
    "SEMESTRES — Périodes",
    "ETUDIANTS — 800 étudiants",
]
tables_col2 = [
    "ENSEIGNANTS — 50 enseignants",
    "MATIERES — 50 matières",
    "AFFECTATIONS — Enseignant↔Matière",
    "SEANCES — 3 000 séances",
]
tables_col3 = [
    "ABSENCES — 120 000+ absences",
    "JUSTIFICATIFS — Documents",
    "AUDIT_LOG — Journal d'audit",
    "NOTIFICATIONS — Alertes",
]

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(4), Inches(0.4),
             "Tables principales", 18, GREEN_MAIN, bold=True)
add_bullet_slide(slide, tables_col1, Inches(0.8), Inches(2), Inches(4), Inches(3), 15, GRAY_DARK)

add_text_box(slide, Inches(4.8), Inches(1.5), Inches(4), Inches(0.4),
             "Tables opérationnelles", 18, GREEN_MAIN, bold=True)
add_bullet_slide(slide, tables_col2, Inches(4.8), Inches(2), Inches(4), Inches(3), 15, GRAY_DARK)

add_text_box(slide, Inches(8.8), Inches(1.5), Inches(4), Inches(0.4),
             "Tables support", 18, GREEN_MAIN, bold=True)
add_bullet_slide(slide, tables_col3, Inches(8.8), Inches(2), Inches(4), Inches(3), 15, GRAY_DARK)

# Contraintes
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(12), Inches(0.4),
             "Contraintes appliquées", 18, GREEN_MAIN, bold=True)
items = [
    "PRIMARY KEY (IDENTITY) — Clés primaires auto-incrémentées",
    "FOREIGN KEY — Relations entre tables (CASCADE, SET NULL)",
    "UNIQUE — CNE, matricule, email, combinaisons uniques",
    "CHECK — Statuts valides, dates cohérentes, seuils logiques",
    "NOT NULL — Champs obligatoires sur toutes les tables critiques",
    "23+ INDEX — FK, fonctionnels (UPPER), dates, statuts",
]
add_bullet_slide(slide, items, Inches(0.8), Inches(4.7), Inches(12), Inches(2.5), 15, GRAY_DARK)

add_footer(slide, page)


# ============================================================
# SLIDE 6: PL/SQL
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "4. PL/SQL : Packages, Fonctions, Procédures", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

# PKG_ABSENCES
box = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5), GRAY_LIGHT)
add_text_box(slide, Inches(1), Inches(1.6), Inches(5.5), Inches(0.4),
             "PKG_ABSENCES", 20, GREEN_DARK, bold=True)
items = [
    "Fonctions :",
    "  GET_NB_ABSENCES() — Total absences",
    "  GET_NB_ABSENCES_NON_JUSTIFIEES()",
    "  GET_TAUX_ABSENCE() — Taux en %",
    "  EST_EN_DEPASSEMENT() — Seuil dépassé ?",
    "  GET_STATUT_ETUDIANT_ABSENCES()",
    "",
    "Procédures :",
    "  MARQUER_ABSENCE() — Avec validation",
    "  MARQUER_ABSENCES_BULK() — En masse (CSV)",
    "  JUSTIFIER_ABSENCE() — Soumettre justificatif",
    "  TRAITER_JUSTIFICATIF() — Accepter/Refuser",
    "  VERIFIER_SEUILS_ETUDIANT() — Auto-exclusion",
    "  CLOTURER_SEMESTRE() — Fin de période",
]
add_bullet_slide(slide, items, Inches(1.2), Inches(2.1), Inches(5.2), Inches(4), 14, GRAY_DARK)

# PKG_STATISTIQUES
box = add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3.5), GRAY_LIGHT)
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5), Inches(0.4),
             "PKG_STATISTIQUES", 20, GREEN_DARK, bold=True)
items = [
    "GET_TOP_ABSENTS() — Classement",
    "GET_STATS_PAR_MATIERE() — Par matière",
    "GET_STATS_PAR_GROUPE() — Par groupe",
    "GET_EVOLUTION_MENSUELLE() — Tendances",
    "GENERER_RAPPORT_ABSENCES() — Rapport",
    "",
    "Retourne SYS_REFCURSOR pour",
    "consommation par Django",
]
add_bullet_slide(slide, items, Inches(7.4), Inches(2.1), Inches(5), Inches(3), 14, GRAY_DARK)

# Transactions
box = add_shape(slide, Inches(7), Inches(5.3), Inches(5.5), Inches(1.3), GRAY_LIGHT)
add_text_box(slide, Inches(7.2), Inches(5.4), Inches(5), Inches(0.4),
             "Gestion des Transactions", 18, GREEN_DARK, bold=True)
items = [
    "SAVEPOINT, COMMIT, ROLLBACK",
    "Exceptions personnalisées (RAISE_APPLICATION_ERROR)",
]
add_bullet_slide(slide, items, Inches(7.4), Inches(5.9), Inches(5), Inches(0.8), 14, GRAY_DARK)

add_footer(slide, page)


# ============================================================
# SLIDE 7: Triggers
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "5. Triggers et Audit Automatique", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

# Audit triggers
box = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.8), GRAY_LIGHT)
add_text_box(slide, Inches(1), Inches(1.6), Inches(5.5), Inches(0.4),
             "Triggers d'Audit (3)", 20, RED, bold=True)
items = [
    "TRG_AUDIT_ABSENCES — Log INSERT/UPDATE/DELETE",
    "TRG_AUDIT_ETUDIANTS — Suivi modifications",
    "TRG_AUDIT_JUSTIFICATIFS — Traçabilité",
    "",
    "→ Stockent anciennes et nouvelles valeurs",
    "→ Table AUDIT_LOG avec CLOB pour les détails",
]
add_bullet_slide(slide, items, Inches(1.2), Inches(2.1), Inches(5.2), Inches(2), 14, GRAY_DARK)

# Business triggers
box = add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(2.8), GRAY_LIGHT)
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5), Inches(0.4),
             "Triggers Métier (5)", 20, GREEN_MAIN, bold=True)
items = [
    "TRG_ABSENCE_DATE_AUTO — Auto-remplir la date",
    "TRG_PROTECT_ABSENCE_JUSTIFIEE — Protection",
    "TRG_CHECK_ETUDIANT_ACTIF — Validation statut",
    "TRG_NOTIF_CHANGEMENT_STATUT — Notification",
    "TRG_SEANCE_EFFECTUEE — Mise à jour auto",
]
add_bullet_slide(slide, items, Inches(7.4), Inches(2.1), Inches(5), Inches(2), 14, GRAY_DARK)

# Screenshot
add_screenshot(slide, "admin_audit.png", Inches(1.5), Inches(4.6), Inches(10), Inches(2.3))

add_footer(slide, page)


# ============================================================
# SLIDE 8: Vues et Vues Matérialisées
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "6. Vues et Vues Matérialisées", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

# Vues
box = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.5), GRAY_LIGHT)
add_text_box(slide, Inches(1), Inches(1.6), Inches(5.5), Inches(0.4),
             "Vues classiques (5)", 20, GREEN_MAIN, bold=True)
items = [
    "V_ABSENCES_DETAIL — Détails complets",
    "V_RESUME_ABSENCES_ETUDIANT — Résumé/étudiant",
    "V_SEANCES_ABSENTS — Séances + nb absents",
    "V_JUSTIFICATIFS_EN_ATTENTE — À traiter",
    "V_NOTIFICATIONS_NON_LUES — Alertes",
]
add_bullet_slide(slide, items, Inches(1.2), Inches(2.1), Inches(5.2), Inches(2), 14, GRAY_DARK)

# Materialized views
box = add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(2.5), GRAY_LIGHT)
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5), Inches(0.4),
             "Vues Matérialisées (4)", 20, ORANGE, bold=True)
items = [
    "MV_STATS_PAR_GROUPE — Dashboard groupes",
    "MV_STATS_PAR_MATIERE — Dashboard matières",
    "MV_EVOLUTION_MENSUELLE — Tendances",
    "MV_KPI_GLOBAL — Indicateurs temps réel",
    "→ REFRESH ON DEMAND via REFRESH_ALL_MV()",
]
add_bullet_slide(slide, items, Inches(7.4), Inches(2.1), Inches(5), Inches(2), 14, GRAY_DARK)

# Screenshot stats
add_screenshot(slide, "admin_stats_groupes.png", Inches(1.5), Inches(4.3), Inches(10), Inches(2.5))

add_footer(slide, page)


# ============================================================
# SLIDE 9: Optimisation
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "7. Optimisation et EXPLAIN PLAN", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

items = [
    "DBMS_STATS.GATHER_TABLE_STATS — Statistiques pour l'optimiseur",
    "",
    "EXPLAIN PLAN FOR — Analyse du plan d'exécution de chaque requête",
    "DBMS_XPLAN.DISPLAY — Affichage formaté du plan",
    "",
    "Index créés pour optimiser :",
    "  • 13 index sur clés étrangères (jointures rapides)",
    "  • Index fonctionnel UPPER(nom, prenom) pour recherche insensible à la casse",
    "  • Index sur date_seance et date_absence (requêtes temporelles)",
    "  • Index sur statut et est_justifiee (filtres fréquents)",
    "",
    "Comparaison mesurée (SET TIMING ON) :",
    "  • Requête directe sur tables : ~200ms",
    "  • Vue matérialisée MV_STATS_PAR_GROUPE : ~5ms (×40 plus rapide)",
    "",
    "Volume de données testé : 120 000+ absences, 3 000 séances, 800 étudiants",
]
add_bullet_slide(slide, items, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5), 16, GRAY_DARK)

add_footer(slide, page)


# ============================================================
# SLIDE 10: Demo App - Dashboard
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "8. Application Django — Tableau de Bord", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

add_screenshot(slide, "admin_dashboard.png", Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))

add_footer(slide, page)


# ============================================================
# SLIDE 11: Demo - CRUD
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "8. Application Django — CRUD et Absences", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

# Two screenshots side by side
add_screenshot(slide, "admin_etudiants.png", Inches(0.5), Inches(1.4), Inches(6.2), Inches(2.7))
add_screenshot(slide, "admin_absences.png", Inches(6.8), Inches(1.4), Inches(6.2), Inches(2.7))

add_text_box(slide, Inches(0.5), Inches(4.2), Inches(6.2), Inches(0.3),
             "Liste des étudiants (recherche, filtres, pagination)", 12, GRAY_MED, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.8), Inches(4.2), Inches(6.2), Inches(0.3),
             "Liste des absences (filtres par groupe, matière, date)", 12, GRAY_MED, alignment=PP_ALIGN.CENTER)

add_screenshot(slide, "admin_etudiant_detail.png", Inches(2), Inches(4.6), Inches(9), Inches(2.2))
add_text_box(slide, Inches(2), Inches(6.85), Inches(9), Inches(0.3),
             "Détail étudiant — KPIs Oracle (fonctions PL/SQL appelées en temps réel)", 12, GRAY_MED, alignment=PP_ALIGN.CENTER)

add_footer(slide, page)


# ============================================================
# SLIDE 12: Roles
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "9. Gestion des Rôles", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

# 3 columns for 3 roles
# Admin
box = add_shape(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(2.2), GRAY_LIGHT)
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(3.6), Inches(0.4),
             "Admin", 20, GREEN_DARK, bold=True, alignment=PP_ALIGN.CENTER)
items = ["Tout accès", "CRUD complet", "Gérer utilisateurs", "Audit + traitements"]
add_bullet_slide(slide, items, Inches(0.9), Inches(2.1), Inches(3.4), Inches(1.5), 14, GRAY_DARK)

# Enseignant
box = add_shape(slide, Inches(4.8), Inches(1.5), Inches(4), Inches(2.2), GRAY_LIGHT)
add_text_box(slide, Inches(5), Inches(1.6), Inches(3.6), Inches(0.4),
             "Enseignant", 20, GREEN_DARK, bold=True, alignment=PP_ALIGN.CENTER)
items = ["Ses séances", "Marquer absences", "Voir ses étudiants", "Statistiques"]
add_bullet_slide(slide, items, Inches(5.2), Inches(2.1), Inches(3.4), Inches(1.5), 14, GRAY_DARK)

# Etudiant
box = add_shape(slide, Inches(9.1), Inches(1.5), Inches(4), Inches(2.2), GRAY_LIGHT)
add_text_box(slide, Inches(9.3), Inches(1.6), Inches(3.6), Inches(0.4),
             "Étudiant", 20, GREEN_DARK, bold=True, alignment=PP_ALIGN.CENTER)
items = ["Ses absences", "Soumettre justificatif", "Ses notifications", "Lecture seule"]
add_bullet_slide(slide, items, Inches(9.5), Inches(2.1), Inches(3.4), Inches(1.5), 14, GRAY_DARK)

# Screenshots
add_screenshot(slide, "enseignant_dashboard.png", Inches(0.3), Inches(4), Inches(6.3), Inches(2.8))
add_screenshot(slide, "etudiant_dashboard.png", Inches(6.8), Inches(4), Inches(6.3), Inches(2.8))
add_text_box(slide, Inches(0.3), Inches(6.85), Inches(6.3), Inches(0.3),
             "Vue Enseignant", 12, GRAY_MED, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.8), Inches(6.85), Inches(6.3), Inches(0.3),
             "Vue Étudiant", 12, GRAY_MED, alignment=PP_ALIGN.CENTER)

add_footer(slide, page)


# ============================================================
# SLIDE 13: Login + Justificatifs
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "8. Connexion et Justificatifs", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

add_screenshot(slide, "page_login.png", Inches(0.5), Inches(1.4), Inches(6), Inches(3))
add_screenshot(slide, "admin_justificatifs.png", Inches(6.8), Inches(1.4), Inches(6), Inches(3))

add_text_box(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(0.3),
             "Page de connexion", 12, GRAY_MED, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.8), Inches(4.5), Inches(6), Inches(0.3),
             "Gestion des justificatifs (Admin)", 12, GRAY_MED, alignment=PP_ALIGN.CENTER)

add_screenshot(slide, "admin_traitements.png", Inches(2), Inches(4.9), Inches(9), Inches(2))
add_text_box(slide, Inches(2), Inches(6.95), Inches(9), Inches(0.3),
             "Traitements Oracle — Appels directs aux procédures PL/SQL", 12, GRAY_MED, alignment=PP_ALIGN.CENTER)

add_footer(slide, page)


# ============================================================
# SLIDE 14: Export & Backup
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "10. Export et Sauvegarde", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

# Export
box = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(3), GRAY_LIGHT)
add_text_box(slide, Inches(1), Inches(1.6), Inches(5.5), Inches(0.4),
             "Export depuis Django", 20, GREEN_MAIN, bold=True)
items = [
    "Export Excel (.xlsx) — openpyxl",
    "  → En-têtes stylés, colonnes ajustées",
    "  → Jusqu'à 5 000 lignes exportées",
    "",
    "Export PDF — ReportLab",
    "  → Rapport formaté avec en-tête ISCAE",
    "  → Tableau paginé automatiquement",
    "  → Jusqu'à 500 lignes par export",
]
add_bullet_slide(slide, items, Inches(1.2), Inches(2.1), Inches(5.2), Inches(2.5), 14, GRAY_DARK)

# Backup
box = add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3), GRAY_LIGHT)
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5), Inches(0.4),
             "Sauvegarde Oracle (Data Pump)", 20, GREEN_MAIN, bold=True)
items = [
    "EXPDP — Export complet du schéma",
    "IMPDP — Import/restauration",
    "DIRECTORY DATA_PUMP_DIR configuré",
    "",
    "Script batch automatisé :",
    "  → Sauvegarde quotidienne à 02h00",
    "  → Compression des dumps",
    "  → Rotation automatique (30 jours)",
]
add_bullet_slide(slide, items, Inches(7.4), Inches(2.1), Inches(5), Inches(2.5), 14, GRAY_DARK)

# Connection architecture
box = add_shape(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.8), GRAY_LIGHT)
add_text_box(slide, Inches(1), Inches(4.9), Inches(11), Inches(0.4),
             "Connexion Django ↔ Oracle", 20, GREEN_MAIN, bold=True)
items = [
    "ENGINE: django.db.backends.oracle  →  oracledb (compatible cx_Oracle)  →  Oracle 18c XE (localhost:1521/XE)",
    "Models avec managed=False — Django lit/écrit sans modifier le schéma Oracle",
    "oracle_services.py — cursor.callfunc() et cursor.callproc() pour appeler le PL/SQL",
]
add_bullet_slide(slide, items, Inches(1.2), Inches(5.4), Inches(11), Inches(1.2), 14, GRAY_DARK)

add_footer(slide, page)


# ============================================================
# SLIDE 15: Conclusion
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), GREEN_MAIN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "11. Conclusion et Perspectives", 32, GREEN_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), GREEN_MAIN)

# Réalisations
box = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(3.5), GRAY_LIGHT)
add_text_box(slide, Inches(1), Inches(1.6), Inches(5.5), Inches(0.4),
             "Ce qui a été réalisé", 20, GREEN_MAIN, bold=True)
items = [
    "Base Oracle complète (12 tables, 120K+ lignes)",
    "PL/SQL avancé (packages, fonctions, procédures)",
    "8 triggers (audit + règles métier)",
    "9 vues (classiques + matérialisées)",
    "Optimisation avec EXPLAIN PLAN",
    "Application Django fonctionnelle",
    "3 rôles utilisateurs distincts",
    "Export PDF/Excel + sauvegarde Data Pump",
]
add_bullet_slide(slide, items, Inches(1.2), Inches(2.1), Inches(5.2), Inches(3), 15, GRAY_DARK)

# Perspectives
box = add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3.5), GRAY_LIGHT)
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5), Inches(0.4),
             "Perspectives d'évolution", 20, ORANGE, bold=True)
items = [
    "Déploiement en production (serveur)",
    "Intégration avec le SI de l'ISCAE",
    "Notifications par email/SMS",
    "Application mobile",
    "Dashboard en temps réel (WebSocket)",
    "Partitionnement Oracle pour le volume",
    "Oracle RAC pour la haute disponibilité",
    "Intégration avec un ERP académique",
]
add_bullet_slide(slide, items, Inches(7.4), Inches(2.1), Inches(5), Inches(3), 15, GRAY_DARK)

add_footer(slide, page)


# ============================================================
# SLIDE 16: Merci
# ============================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Top bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), GREEN_MAIN)

# Logo
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(5.6), Inches(1.2), Inches(2), Inches(2))

add_text_box(slide, Inches(1), Inches(3.5), Inches(11.333), Inches(1),
             "Merci pour votre attention", 40, GREEN_DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(4.5), Inches(2.333), Inches(0.04), GREEN_MAIN)

add_text_box(slide, Inches(1), Inches(4.8), Inches(11.333), Inches(0.6),
             "Des questions ?", 28, GRAY_MED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11.333), Inches(0.5),
             "IP22240 — Cherifa Mohamed Yeslem Bellahi", 16, GRAY_DARK, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5),
             "IP22237 — El Vagih Ahmed Zeine", 16, GRAY_DARK, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), GREEN_MAIN)
add_text_box(slide, Inches(1), Inches(7.12), Inches(11.333), Inches(0.35),
             "ISCAE Nouakchott — Master MPIAG — 2025-2026", 11, WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {page}")
